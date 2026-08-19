from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_FILE = "景区导览AI数字人系统技术汇报.pptx"
ROOT = Path(__file__).resolve().parent

BG = RGBColor(245, 247, 251)
NAVY = RGBColor(15, 23, 42)
BLUE = RGBColor(37, 99, 235)
SKY = RGBColor(96, 165, 250)
MINT = RGBColor(16, 185, 129)
AMBER = RGBColor(245, 158, 11)
RED = RGBColor(239, 68, 68)
SLATE = RGBColor(71, 85, 105)
LIGHT = RGBColor(226, 232, 240)
WHITE = RGBColor(255, 255, 255)


@dataclass
class BulletBlock:
    title: str
    items: list[str]
    color: RGBColor = BLUE


def set_bg(slide) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.22), Inches(12.0), Inches(0.65))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.name = "Microsoft YaHei"
    run.font.color.rgb = NAVY
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.48), Inches(0.78), Inches(8.2), Inches(0.35))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(10.5)
        r2.font.name = "Microsoft YaHei"
        r2.font.color.rgb = SLATE
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.45), Inches(1.08), Inches(12.0), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()


def add_footer(slide, index_text: str) -> None:
    box = slide.shapes.add_textbox(Inches(11.35), Inches(6.95), Inches(1.0), Inches(0.25))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = index_text
    r.font.size = Pt(9)
    r.font.name = "Microsoft YaHei"
    r.font.color.rgb = SLATE


def add_card(slide, left, top, width, height, title: str, body_lines: list[str], accent: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LIGHT

    tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, Inches(0.14), height)
    tag.fill.solid()
    tag.fill.fore_color.rgb = accent
    tag.line.fill.background()

    title_box = slide.shapes.add_textbox(left + Inches(0.24), top + Inches(0.12), width - Inches(0.34), Inches(0.35))
    tf1 = title_box.text_frame
    p1 = tf1.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.size = Pt(14)
    r1.font.bold = True
    r1.font.name = "Microsoft YaHei"
    r1.font.color.rgb = NAVY

    body = slide.shapes.add_textbox(left + Inches(0.24), top + Inches(0.52), width - Inches(0.34), height - Inches(0.62))
    tf2 = body.text_frame
    tf2.word_wrap = True
    tf2.margin_left = 0
    tf2.margin_right = 0
    for idx, line in enumerate(body_lines):
        p = tf2.paragraphs[0] if idx == 0 else tf2.add_paragraph()
        p.text = f"• {line}"
        p.level = 0
        p.space_after = Pt(4)
        for run in p.runs:
            run.font.size = Pt(11)
            run.font.name = "Microsoft YaHei"
            run.font.color.rgb = SLATE


def add_section_bullets(slide, left, top, width, height, block: BulletBlock) -> None:
    add_card(slide, left, top, width, height, block.title, block.items, block.color)


def add_flow_box(slide, left, top, width, height, text: str, fill: RGBColor, font_size: int = 12) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.bold = True
    r.font.size = Pt(font_size)
    r.font.name = "Microsoft YaHei"
    r.font.color.rgb = WHITE if fill != LIGHT else NAVY


def add_arrow(slide, left, top, width, height, color: RGBColor = SKY) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color


def add_small_note(slide, text: str, left, top, width) -> None:
    box = slide.shapes.add_textbox(left, top, width, Inches(0.28))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(10)
    r.font.name = "Microsoft YaHei"
    r.font.color.rgb = SLATE


def add_table_like(slide, left, top, width, headers: list[str], rows: list[list[str]], col_widths: list[float]) -> None:
    row_count = len(rows) + 1
    table = slide.shapes.add_table(row_count, len(headers), left, top, width, Inches(0.35 + row_count * 0.34)).table
    for idx, head in enumerate(headers):
        cell = table.cell(0, idx)
        cell.text = head
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(10.5)
                r.font.bold = True
                r.font.name = "Microsoft YaHei"
                r.font.color.rgb = WHITE
    for c, w in enumerate(col_widths):
        table.columns[c].width = Inches(w)
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(9.5)
                    run.font.name = "Microsoft YaHei"
                    run.font.color.rgb = NAVY if c_idx == 0 else SLATE


def title_slide(prs: Presentation, idx: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    banner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.75), Inches(11.7), Inches(5.6))
    banner.fill.solid()
    banner.fill.fore_color.rgb = WHITE
    banner.line.color.rgb = LIGHT

    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(0.75), Inches(0.22), Inches(5.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.95), Inches(1.18), Inches(10.6), Inches(1.4))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "景区导览 AI 数字人系统\n整体分析与技术汇报"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.name = "Microsoft YaHei"
    r.font.color.rgb = NAVY

    info = [
        "项目对象：游客前台 + 管理后台 + FastAPI 后端 + RAG 知识库 + ASR/TTS + 数字人视频能力",
        "汇报重点：模块划分、关键调用链、核心实现细节、部署方式、当前能力边界与优化方向",
        "分析结论：系统已具备从资料入库到游客问答再到运营分析的完整闭环能力",
    ]
    add_card(slide, Inches(0.95), Inches(2.6), Inches(10.75), Inches(2.35), "汇报定位", info, MINT)

    add_small_note(slide, "代码仓库结构：web-client / admin-web / backend / docker-compose / 资料包", Inches(0.95), Inches(5.2), Inches(9.8))
    add_footer(slide, f"{idx:02d}")


def overview_slide(prs: Presentation, idx: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "1. 系统总体概览", "该项目是一个面向景区导览场景的 AI 数字人单体系统")
    blocks = [
        BulletBlock("核心定位", ["将官方景区资料构造成知识库，为游客提供文本/语音问答、数字人讲解和路线推荐"], BLUE),
        BulletBlock("前台能力", ["Vue3 + Pinia 单页应用", "支持文本输入、按住说话、数字人讲解、路线推荐弹窗"], MINT),
        BulletBlock("后台能力", ["React + Ant Design 管理控制台", "提供数据大屏、知识库上传、感受度分析"], AMBER),
        BulletBlock("后端能力", ["FastAPI 单体服务", "负责 RAG、LLM、ASR、TTS、日志分析、数字人视频代理"], RED),
    ]
    positions = [
        (Inches(0.55), Inches(1.45)),
        (Inches(6.45), Inches(1.45)),
        (Inches(0.55), Inches(4.05)),
        (Inches(6.45), Inches(4.05)),
    ]
    for block, (left, top) in zip(blocks, positions):
        add_section_bullets(slide, left, top, Inches(5.35), Inches(2.2), block)
    add_footer(slide, f"{idx:02d}")


def architecture_slide(prs: Presentation, idx: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "2. 分层架构与组件关系", "从浏览器接入到 AI 能力层和数据层形成清晰的单体分层")

    add_flow_box(slide, Inches(0.65), Inches(1.45), Inches(2.2), Inches(0.7), "游客端\nVue3 / Pinia / Vite", BLUE)
    add_flow_box(slide, Inches(0.65), Inches(2.45), Inches(2.2), Inches(0.7), "管理后台\nReact / AntD / ECharts", SKY)
    add_flow_box(slide, Inches(3.25), Inches(1.95), Inches(2.0), Inches(0.8), "Nginx\n静态资源 + /api 反代", NAVY)
    add_flow_box(slide, Inches(5.75), Inches(1.95), Inches(2.5), Inches(0.8), "FastAPI Backend\nAPI / 编排 / 鉴权 / 静态音频", MINT)
    add_flow_box(slide, Inches(8.75), Inches(1.2), Inches(2.8), Inches(0.7), "RAG 与知识处理\nChroma + bge-m3 + 文档解析", AMBER)
    add_flow_box(slide, Inches(8.75), Inches(2.15), Inches(2.8), Inches(0.7), "语音与数字人\nfaster-whisper + Edge-TTS + 讯飞视频", RED)
    add_flow_box(slide, Inches(8.75), Inches(3.1), Inches(2.8), Inches(0.7), "会话与分析\nSQLite / Redis / 日志统计", BLUE)

    add_arrow(slide, Inches(2.9), Inches(1.67), Inches(0.28), Inches(0.25))
    add_arrow(slide, Inches(2.9), Inches(2.67), Inches(0.28), Inches(0.25))
    add_arrow(slide, Inches(5.35), Inches(2.17), Inches(0.28), Inches(0.25))
    add_arrow(slide, Inches(8.35), Inches(1.42), Inches(0.28), Inches(0.25))
    add_arrow(slide, Inches(8.35), Inches(2.37), Inches(0.28), Inches(0.25))
    add_arrow(slide, Inches(8.35), Inches(3.32), Inches(0.28), Inches(0.25))

    add_card(
        slide,
        Inches(0.65),
        Inches(4.35),
        Inches(11.0),
        Inches(1.9),
        "架构判断",
        [
            "这是典型的 AI 业务单体架构，而不是微服务；优点是开发集成快，适合竞赛和演示场景。",
            "Nginx 将两个前端和后端 API 聚合到同一入口，降低跨域与部署复杂度。",
            "核心价值不在 CRUD，而在“知识库检索 + 语音链路 + 数字人呈现 + 运营闭环”的组合能力。",
        ],
        BLUE,
    )
    add_footer(slide, f"{idx:02d}")


def deploy_slide(prs: Presentation, idx: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "3. 部署拓扑与运行环境", "Docker Compose 提供 Redis、Backend、Nginx 的一键编排")
    add_table_like(
        slide,
        Inches(0.65),
        Inches(1.45),
        Inches(11.0),
        ["服务", "职责", "端口/形态", "关键依赖"],
        [
            ["nginx", "承载游客端与管理后台静态资源，并反向代理 /api 与 /static", "80", "docker/nginx/Dockerfile"],
            ["backend", "FastAPI + uvicorn，承载聊天、ASR、TTS、RAG、管理接口", "内部 8000", "Python 3.10 / requirements.txt"],
            ["redis", "缓存会话记忆；未连接时可降级为进程内内存字典", "内部 6379", "redis:7-alpine"],
            ["持久化卷", "保存 SQLite、Chroma、TTS 资源、HF 模型缓存", "宿主卷", "backend_data / redis_data / hf_cache"],
        ],
        [1.2, 4.6, 1.4, 3.2],
    )
    add_card(
        slide,
        Inches(0.65),
        Inches(4.15),
        Inches(5.3),
        Inches(1.85),
        "启动过程",
        [
            "启动 backend 时自动执行资料同步与向量库初始化。",
            "首次启动需下载 bge-m3 模型，因此冷启动时间较长。",
            "所有前端默认通过 Nginx 与 Vite 代理访问 /api。",
        ],
        MINT,
    )
    add_card(
        slide,
        Inches(6.15),
        Inches(4.15),
        Inches(5.5),
        Inches(1.85),
        "运维侧特征",
        [
            "数据层以 SQLite + Chroma 为主，易部署但不适合大规模高并发。",
            "配置项集中在 backend/.env，支持切换 LLM、ASR、TTS、鉴权和索引策略。",
            "整个系统更偏单机演示架构，而不是生产级高可用集群。",
        ],
        AMBER,
    )
    add_footer(slide, f"{idx:02d}")


def backend_slide(prs: Presentation, idx: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "4. 后端目录与模块职责", "backend/app 按 API、services、models、schemas、core、utils 分层")
    blocks = [
        BulletBlock("api", ["v1 下按 chat / asr / tts / recommend / avatar / admin 划分接口入口"], BLUE),
        BulletBlock("services", ["承载业务编排与 AI 能力接入，是真正的核心层"], MINT),
        BulletBlock("models & schemas", ["SQLAlchemy 模型保存文档和会话；Pydantic 模型定义接口契约"], AMBER),
        BulletBlock("core & utils", ["配置、数据库、鉴权、统一响应、ID 生成、文本分块等基础能力"], RED),
    ]
    for i, block in enumerate(blocks):
        left = Inches(0.65 + (i % 2) * 5.55)
        top = Inches(1.5 + (i // 2) * 2.2)
        add_section_bullets(slide, left, top, Inches(5.1), Inches(1.8), block)
    add_card(
        slide,
        Inches(0.65),
        Inches(5.95),
        Inches(11.0),
        Inches(0.8),
        "结论",
        [
            "后端设计并不复杂，但分层比较工整，适合在汇报中强调“职责清晰、便于替换模型与前端组件”。"
        ],
        BLUE,
    )
    add_footer(slide, f"{idx:02d}")


def chat_chain_slide(prs: Presentation, idx: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "5. 核心问答链路", "聊天主流程由 ChatService 串联：记忆 -> 检索 -> 生成 -> 合成 -> 返回")

    steps = [
        ("1", "读取会话记忆", "从 sessionId 获取最近 3 轮历史，保证上下文追问成立", BLUE),
        ("2", "RAG 检索", "通过 Chroma 检索 top-3 文档块并抽取 source 列表", MINT),
        ("3", "LLM 生成", "把参考资料、历史对话、兴趣偏好拼成 Prompt 调用模型", AMBER),
        ("4", "TTS 合成", "将 answerText 合成为音频与 phonemes / duration", RED),
        ("5", "会话更新", "把本轮用户问题与回答写回会话记忆，形成闭环", BLUE),
    ]
    for i, (num, title, desc, color) in enumerate(steps):
        top = Inches(1.55 + i * 1.0)
        add_flow_box(slide, Inches(0.8), top, Inches(0.6), Inches(0.56), num, color)
        add_card(slide, Inches(1.55), top - Inches(0.06), Inches(10.0), Inches(0.72), title, [desc], color)

    add_small_note(
        slide,
        "置信度是检索置信度与 LLM 置信度的组合值；若无检索命中，系统会降到低置信度并返回礼貌兜底话术。",
        Inches(0.8),
        Inches(6.7),
        Inches(10.9),
    )
    add_footer(slide, f"{idx:02d}")


def rag_slide(prs: Presentation, idx: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "6. 知识库处理与 RAG 技术细节", "系统依赖官方资料包构建本地知识检索能力")
    add_card(
        slide,
        Inches(0.65),
        Inches(1.45),
        Inches(5.4),
        Inches(4.7),
        "文档处理流水线",
        [
            "支持 .pdf / .txt / .md / .docx / .xlsx 五类输入，其中 xlsx 会被转成可检索的结构化景点文本。",
            "处理状态流转为 pending -> parsing -> chunking -> indexing -> ready/failed，便于后台轮询展示进度。",
            "分块参数由配置驱动，默认 chunk_size=500、overlap=50。",
            "重复上传时会先删除旧向量，再按批次写入 Chroma，避免索引污染。",
        ],
        AMBER,
    )
    add_card(
        slide,
        Inches(6.25),
        Inches(1.45),
        Inches(5.4),
        Inches(4.7),
        "向量检索实现",
        [
            "Chroma 使用 PersistentClient 落盘，便于单机部署与演示。",
            "嵌入模型默认采用 BAAI/bge-m3，通过 sentence-transformers 延迟加载。",
            "相似度空间配置为 cosine，并将距离映射为 [0,1] 置信度供上层合成。",
            "这套设计很适合竞赛，但在大知识库场景中需要进一步考虑索引规模与查询时延。",
        ],
        BLUE,
    )
    add_footer(slide, f"{idx:02d}")


def ai_slide(prs: Presentation, idx: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "7. LLM / ASR / TTS / 数字人链路", "系统把语义理解、语音处理和呈现层结合成完整的人机交互链")
    blocks = [
        BulletBlock("LLM 层", ["支持 mock 与 OpenAI 兼容接口两种模式，可接 DeepSeek 等平台", "Prompt 明确要求严格依据官方资料，降低幻觉"], BLUE),
        BulletBlock("ASR 层", ["后端支持 faster-whisper 与 OpenAI Whisper 路径", "前端录音会重采样为 16kHz 单声道 WAV，贴合识别需求"], MINT),
        BulletBlock("TTS 层", ["Edge-TTS 负责语音合成，返回 audioUrl、durationMs、phonemes", "为字幕节奏和口型同步提供结构化时间信息"], AMBER),
        BulletBlock("数字人层", ["当前首页实际使用 XfyunAvatar 组件，走“文本 -> 本地播报占位 -> 服务端生成同步视频”", "代码中还保留一套 Live2D + phoneme 口型同步备用链路"], RED),
    ]
    positions = [
        (Inches(0.65), Inches(1.5)),
        (Inches(6.2), Inches(1.5)),
        (Inches(0.65), Inches(4.05)),
        (Inches(6.2), Inches(4.05)),
    ]
    for block, (left, top) in zip(blocks, positions):
        add_section_bullets(slide, left, top, Inches(5.1), Inches(2.1), block)
    add_footer(slide, f"{idx:02d}")


def frontend_slide(prs: Presentation, idx: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "8. 游客前台模块拆解", "web-client 是单页应用，重点在交互体验与数字人承载")
    add_card(
        slide,
        Inches(0.65),
        Inches(1.45),
        Inches(3.5),
        Inches(4.95),
        "页面组织",
        [
            "不使用 vue-router，入口由 App.vue 直接拼接头部、数字人区和聊天区。",
            "ChatInterface 承担欢迎态、消息流、快捷问题、路线弹窗和输入区域。",
            "RouteRecommendModal 通过 Teleport 方式出现，支持列表视图与 Leaflet 地图。",
        ],
        BLUE,
    )
    add_card(
        slide,
        Inches(4.45),
        Inches(1.45),
        Inches(3.55),
        Inches(4.95),
        "状态管理",
        [
            "全局仅有 chat Store，保存 sessionId、messages、lastAnswer、lastEmotion、lastTts 等信息。",
            "sendText 与 sendVoice 最终复用同一问答接口，体现了前端状态建模的收敛性。",
            "聊天区会根据 durationMs 驱动打字机效果，让文字揭示速度与语音节奏更一致。",
        ],
        MINT,
    )
    add_card(
        slide,
        Inches(8.25),
        Inches(1.45),
        Inches(3.4),
        Inches(4.95),
        "语音与数字人",
        [
            "按住说话组件通过 Web Audio 采 PCM，再编码为 WAV 上传 ASR。",
            "当前主展示链路是 XfyunAvatar，而不是 Live2D；这说明项目优先保障“可展示视频效果”。",
            "代码中存在更细粒度的 phoneme 口型同步实现，具备继续演进空间。",
        ],
        RED,
    )
    add_footer(slide, f"{idx:02d}")


def admin_slide(prs: Presentation, idx: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "9. 管理后台模块拆解", "admin-web 负责知识维护、运营观测和情感分析")
    add_table_like(
        slide,
        Inches(0.65),
        Inches(1.45),
        Inches(11.0),
        ["页面", "核心能力", "使用接口", "技术说明"],
        [
            ["Dashboard", "展示会话数、消息数、热问答、满意度趋势", "/api/v1/admin/dashboard/overview", "30 秒轮询刷新，图表基于 ECharts"],
            ["KnowledgeManage", "上传资料、轮询索引状态、查看文档处理结果", "/api/v1/admin/kb/documents", "支持 PDF/TXT/MD/DOCX 上传"],
            ["SentimentAnalysis", "展示近 7 日情感占比与热点话题词云", "/api/v1/admin/analytics/sentiment-trend", "使用 echarts-wordcloud 做可视化"],
        ],
        [1.6, 3.1, 2.8, 3.5],
    )
    add_card(
        slide,
        Inches(0.65),
        Inches(4.45),
        Inches(5.3),
        Inches(1.7),
        "权限模型",
        [
            "前端默认使用本地 token，缺失时回退到硬编码管理员 token。",
            "后端支持 JWT 或固定 ADMIN_API_TOKEN，当前不是细粒度 RBAC。",
        ],
        AMBER,
    )
    add_card(
        slide,
        Inches(6.15),
        Inches(4.45),
        Inches(5.5),
        Inches(1.7),
        "系统定位",
        [
            "后台更偏演示型运营后台，而非企业级多角色内容平台。",
            "它的价值在于支撑 AI 数字人的资料入库和运营分析，而不是复杂业务审批流。",
        ],
        BLUE,
    )
    add_footer(slide, f"{idx:02d}")


def data_slide(prs: Presentation, idx: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "10. 数据模型与持久化设计", "结构化库负责元数据与日志，向量库负责语义检索")
    add_table_like(
        slide,
        Inches(0.65),
        Inches(1.45),
        Inches(11.0),
        ["数据对象", "关键字段", "用途", "存储位置"],
        [
            ["KnowledgeDocument", "doc_id / title / category / status / progress / chunk_count", "追踪知识文档入库进度", "SQLite"],
            ["ChatSession", "session_id / preference / created_at / updated_at", "记录会话级信息", "SQLite"],
            ["ChatMessage", "session_id / user_text / assistant_text / sentiment / confidence / latency", "支撑日志分析与热点统计", "SQLite"],
            ["向量块元数据", "doc_id / title / source_file / chunk_index / scenic_area_id", "支撑 Chroma 检索与来源展示", "Chroma"],
            ["会话缓存", "最近轮次历史对话", "提高多轮问答读取效率", "Redis 或内存"],
        ],
        [1.9, 3.9, 2.4, 1.8],
    )
    add_small_note(slide, "当前方案对演示非常友好；若面向真实景区长期运营，SQLite 与本地向量库会成为扩展瓶颈。", Inches(0.65), Inches(6.5), Inches(10.9))
    add_footer(slide, f"{idx:02d}")


def api_slide(prs: Presentation, idx: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "11. 主要 API 模块与交互契约", "接口设计围绕游客主链路与后台管理两条线展开")
    add_table_like(
        slide,
        Inches(0.65),
        Inches(1.45),
        Inches(11.0),
        ["模块", "代表接口", "输入输出重点", "业务意义"],
        [
            ["问答", "POST /api/v1/chat/ask", "输入 sessionId/text/preference，输出 answerText/sources/tts/avatar", "系统主价值接口"],
            ["ASR", "POST /api/v1/asr/recognize", "上传 WAV，返回文本、置信度、时长", "语音输入入口"],
            ["TTS", "POST /api/v1/tts/synthesize", "输出音频地址、时长、音素信息", "数字人音频驱动"],
            ["推荐", "GET|POST /api/v1/recommend/routes", "输出路线列表、兴趣标签、地图要素", "个性化导览"],
            ["后台", "/api/v1/admin/*", "上传知识、查看大屏、分析情感", "运营闭环支撑"],
            ["数字人", "/api/v1/avatar/*", "拉配置、生成视频、转发代理", "呈现层联动"],
        ],
        [1.2, 2.8, 3.9, 3.1],
    )
    add_footer(slide, f"{idx:02d}")


def analytics_slide(prs: Presentation, idx: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "12. 运营分析与价值闭环", "系统不仅能回答问题，还能把交互沉淀为可运营的数据资产")
    add_card(
        slide,
        Inches(0.65),
        Inches(1.45),
        Inches(5.35),
        Inches(4.7),
        "统计逻辑",
        [
            "AnalyticsService 会聚合今日会话数、消息数、热门问题、情感趋势和热点词云。",
            "词云不是简单分词，而是结合中文 2~4 字短语与景区领域关键词加权。",
            "无真实数据时会生成保底结果，保证演示页面不空白。",
        ],
        MINT,
    )
    add_card(
        slide,
        Inches(6.2),
        Inches(1.45),
        Inches(5.45),
        Inches(4.7),
        "业务意义",
        [
            "帮助景区发现游客最关注的问题，如门票、路线、开放时间、历史文化等。",
            "通过情感占比可判断讲解体验与服务满意度走势。",
            "后台由此形成“交互 -> 数据沉淀 -> 运营反馈 -> 知识更新”的闭环。",
        ],
        BLUE,
    )
    add_footer(slide, f"{idx:02d}")


def compare_slide(prs: Presentation, idx: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "13. 设计规格与当前实现对照", "项目既有功能规格说明书，也有已经落地的当前版本代码")
    add_table_like(
        slide,
        Inches(0.65),
        Inches(1.45),
        Inches(11.0),
        ["能力项", "规格设想", "当前实现状态", "评价"],
        [
            ["数字人展示", "Live2D/口型同步/表情切换", "已落地讯飞视频链路，同时保留 Live2D 备用实现", "展示效果优先，工程统一性一般"],
            ["多模态输入", "文本 + 语音 + 流式 ASR", "文本与按住说话已实现，流式 ASR 尚未落地", "MVP 完成度较高"],
            ["RAG 问答", "基于官方资料包的准确问答", "文档解析、Chroma 检索、Prompt 构造已形成闭环", "系统核心亮点"],
            ["推荐系统", "兴趣偏好 + 路线推荐 + 地图可视化", "已实现路线推荐接口和前端弹窗地图", "功能完整但算法偏规则化"],
            ["后台分析", "热点、情感、建议、导出", "热点和情感已实现，导出和建议未完全落地", "演示性强于企业化"],
        ],
        [1.5, 2.7, 4.2, 2.6],
    )
    add_footer(slide, f"{idx:02d}")


def risk_slide(prs: Presentation, idx: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "14. 当前系统优势与风险点", "汇报中既要讲能力，也要讲边界，体现分析深度")
    add_card(
        slide,
        Inches(0.65),
        Inches(1.45),
        Inches(5.35),
        Inches(4.9),
        "优势",
        [
            "端到端链路完整：资料导入、语音问答、数字人播放、后台统计全部打通。",
            "技术选型务实：FastAPI + Chroma + Edge-TTS + Vue/React，便于快速集成。",
            "结构清晰：API、服务层、模型层职责明确，可替换不同 LLM 与数字人实现。",
            "演示友好：默认有 mock 路径与静态兜底，降低现场失败概率。",
        ],
        MINT,
    )
    add_card(
        slide,
        Inches(6.2),
        Inches(1.45),
        Inches(5.45),
        Inches(4.9),
        "风险与改进方向",
        [
            "SQLite + 单机 Chroma 扩展性有限，不适合真实高并发景区运营。",
            "管理员鉴权较轻，缺少细粒度权限、审计日志与正式登录流程。",
            "数字人存在两套实现，当前呈现链路与代码抽象尚未完全统一。",
            "推荐与分析逻辑偏规则化，仍有继续引入更强模型与真实指标体系的空间。",
        ],
        RED,
    )
    add_footer(slide, f"{idx:02d}")


def summary_slide(prs: Presentation, idx: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "15. 汇报总结", "一句话概括：这是一个围绕景区导览场景构建的 AI 数字人闭环系统")
    add_card(
        slide,
        Inches(0.9),
        Inches(1.5),
        Inches(10.5),
        Inches(3.9),
        "最终结论",
        [
            "系统真正的核心不是单个模型，而是“知识检索 + 多轮问答 + 语音合成 + 数字人呈现 + 运营反馈”的组合编排能力。",
            "当前版本已经达到竞赛/答辩可汇报的完整度，尤其适合围绕技术闭环、业务价值和可扩展性展开讲解。",
            "如果后续继续优化，优先级建议依次为：统一数字人链路、增强权限与安全、升级存储与检索架构、提升推荐与分析智能度。",
        ],
        BLUE,
    )
    end_box = slide.shapes.add_textbox(Inches(1.0), Inches(5.85), Inches(10.0), Inches(0.6))
    tf = end_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "谢谢。该 PPT 可直接用于系统汇报，也可继续增补截图、时序图和现场演示页。"
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.name = "Microsoft YaHei"
    r.font.color.rgb = NAVY
    add_footer(slide, f"{idx:02d}")


def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = [
        title_slide,
        overview_slide,
        architecture_slide,
        deploy_slide,
        backend_slide,
        chat_chain_slide,
        rag_slide,
        ai_slide,
        frontend_slide,
        admin_slide,
        data_slide,
        api_slide,
        analytics_slide,
        compare_slide,
        risk_slide,
        summary_slide,
    ]
    for idx, builder in enumerate(slides, start=1):
        builder(prs, idx)

    out = ROOT / OUTPUT_FILE
    prs.save(out)
    return out


if __name__ == "__main__":
    path = build_presentation()
    print(path)
